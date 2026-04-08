from __future__ import annotations

import io
import re
from collections import Counter
from pathlib import Path
from typing import Dict, List

import pandas as pd
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation

from config import settings


KEYWORDS = {
    "sales": ["sales", "units sold", "orders"],
    "revenue": ["revenue", "income", "turnover"],
    "profit": ["profit", "margin", "ebitda"],
    "marketing": ["campaign", "marketing", "lead", "conversion"],
    "finance": ["balance sheet", "cashflow", "expense", "budget"],
}


def _extract_keywords(text: str) -> List[str]:
    lower = text.lower()
    found = []
    for key, terms in KEYWORDS.items():
        if any(t in lower for t in terms):
            found.append(key)
    return sorted(set(found))


def _detect_topic(keywords: List[str]) -> str:
    if not keywords:
        return "general"
    if "sales" in keywords and "revenue" in keywords:
        return "sales_performance"
    if "finance" in keywords or "profit" in keywords:
        return "finance"
    if "marketing" in keywords:
        return "marketing"
    return keywords[0]


def _basic_summary(text: str, max_sentences: int = 3) -> str:
    sentences = re.split(r"(?<=[.!?])\s+", text.strip())
    top = [s for s in sentences if s][:max_sentences]
    if not top:
        return "No summary available."
    return " ".join(top)[:700]


def analyze_tabular_file(file_bytes: bytes, extension: str) -> Dict:
    if extension == ".csv":
        df = pd.read_csv(io.BytesIO(file_bytes))
    else:
        df = pd.read_excel(io.BytesIO(file_bytes))
    sample = df.head(5).to_dict(orient="records")
    text_blob = " ".join(df.columns.astype(str).tolist()) + " " + df.head(20).to_string()
    keywords = _extract_keywords(text_blob)
    summary = (
        f"Dataset with {len(df)} rows and {len(df.columns)} columns. "
        f"Columns include: {', '.join(map(str, df.columns[:8]))}."
    )
    return {
        "content_type": "tabular",
        "columns": [str(c) for c in df.columns],
        "num_rows": int(len(df)),
        "sample_data": sample,
        "text_content": text_blob[: settings.max_text_chars],
        "keywords": keywords,
        "topic": _detect_topic(keywords),
        "dataset_type": _detect_topic(keywords),
        "summary": summary,
    }


def analyze_pdf(file_bytes: bytes) -> Dict:
    reader = PdfReader(io.BytesIO(file_bytes))
    text = "\n".join([(page.extract_text() or "") for page in reader.pages])
    keywords = _extract_keywords(text)
    return {
        "content_type": "document",
        "columns": [],
        "num_rows": 0,
        "sample_data": [],
        "text_content": text[: settings.max_text_chars],
        "keywords": keywords,
        "topic": _detect_topic(keywords),
        "dataset_type": "document",
        "summary": _basic_summary(text),
    }


def analyze_docx(file_bytes: bytes) -> Dict:
    doc = Document(io.BytesIO(file_bytes))
    text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    keywords = _extract_keywords(text)
    return {
        "content_type": "document",
        "columns": [],
        "num_rows": 0,
        "sample_data": [],
        "text_content": text[: settings.max_text_chars],
        "keywords": keywords,
        "topic": _detect_topic(keywords),
        "dataset_type": "document",
        "summary": _basic_summary(text),
    }


def analyze_pptx(file_bytes: bytes) -> Dict:
    prs = Presentation(io.BytesIO(file_bytes))
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
    text = "\n".join(chunks)
    keywords = _extract_keywords(text)
    return {
        "content_type": "presentation",
        "columns": [],
        "num_rows": 0,
        "sample_data": [],
        "text_content": text[: settings.max_text_chars],
        "keywords": keywords,
        "topic": _detect_topic(keywords),
        "dataset_type": "presentation",
        "summary": _basic_summary(text),
    }


def analyze_plain_text(file_bytes: bytes) -> Dict:
    text = file_bytes.decode("utf-8", errors="ignore")
    keywords = _extract_keywords(text)
    words = re.findall(r"[a-zA-Z]{3,}", text.lower())
    common = [w for w, _ in Counter(words).most_common(10)]
    return {
        "content_type": "text",
        "columns": [],
        "num_rows": 0,
        "sample_data": common,
        "text_content": text[: settings.max_text_chars],
        "keywords": keywords,
        "topic": _detect_topic(keywords),
        "dataset_type": "text",
        "summary": _basic_summary(text),
    }


def analyze_file(file_name: str, mime_type: str, file_bytes: bytes) -> Dict:
    ext = Path(file_name).suffix.lower()
    if ext in [".csv", ".xlsx", ".xls"]:
        return analyze_tabular_file(file_bytes, ".csv" if ext == ".csv" else ".xlsx")
    if ext == ".pdf" or mime_type == "application/pdf":
        return analyze_pdf(file_bytes)
    if ext == ".docx" or mime_type.endswith("wordprocessingml.document"):
        return analyze_docx(file_bytes)
    if ext == ".pptx" or mime_type.endswith("presentationml.presentation"):
        return analyze_pptx(file_bytes)
    return analyze_plain_text(file_bytes)
